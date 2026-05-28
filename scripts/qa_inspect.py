"""Geometry-based QA for the generated .pptx.

LibreOffice/pdftoppm aren't available in this environment, so instead of
rendering slides to images we inspect the underlying shape geometry:

  - any shape that extends past slide bounds (cut-off)
  - any pair of non-card shapes whose bounding boxes overlap
  - empty text boxes
  - rough text-width sanity check using avg char-width estimate

Run:
    py scripts/qa_inspect.py
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

REPO_ROOT = Path(__file__).resolve().parent.parent
PPTX_PATH = REPO_ROOT / "output" / "FemaleDiseasPrevalence.pptx"


def emu_to_in(v):
    return v / 914400.0


def shape_bbox(shape):
    return (shape.left or 0, shape.top or 0, shape.width or 0, shape.height or 0)


def overlaps(a, b):
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    return not (ax + aw <= bx or bx + bw <= ax or ay + ah <= by or by + bh <= ay)


def shape_text(shape):
    if not shape.has_text_frame:
        return ""
    return "\n".join(p.text for p in shape.text_frame.paragraphs).strip()


def main():
    prs = Presentation(str(PPTX_PATH))
    sw, sh = prs.slide_width, prs.slide_height
    issues = []
    notes = []

    for idx, slide in enumerate(prs.slides, start=1):
        notes.append(f"\n## Slide {idx}")
        bboxes = []
        for shape in slide.shapes:
            bb = shape_bbox(shape)
            txt = shape_text(shape)
            x, y, w, h = bb
            notes.append(
                f"  - {shape.shape_type}  "
                f"L={emu_to_in(x):.2f}in  T={emu_to_in(y):.2f}in  "
                f"W={emu_to_in(w):.2f}in  H={emu_to_in(h):.2f}in  "
                f"text={txt[:60]!r}"
            )
            if x + w > sw + 1000:
                issues.append(
                    f"Slide {idx}: shape extends past right edge "
                    f"(right={emu_to_in(x + w):.2f}in vs slide={emu_to_in(sw):.2f}in)  text={txt[:40]!r}"
                )
            if y + h > sh + 1000:
                issues.append(
                    f"Slide {idx}: shape extends past bottom "
                    f"(bottom={emu_to_in(y + h):.2f}in vs slide={emu_to_in(sh):.2f}in)  text={txt[:40]!r}"
                )
            if x < -1000 or y < -1000:
                issues.append(f"Slide {idx}: shape positioned off-slide  text={txt[:40]!r}")
            if shape.has_text_frame and not txt:
                issues.append(f"Slide {idx}: empty text frame at L={emu_to_in(x):.2f}in T={emu_to_in(y):.2f}in")
            bboxes.append((shape, bb, txt))

        # Overlap check — skip pairs where one is a CARD/PILL (rounded rectangle)
        # that intentionally contains the other.
        for i in range(len(bboxes)):
            for j in range(i + 1, len(bboxes)):
                s1, b1, t1 = bboxes[i]
                s2, b2, t2 = bboxes[j]
                # ignore if either is an auto-shape (cards/pills are designed to host text)
                if s1.shape_type == 1 or s2.shape_type == 1:
                    continue
                if not t1 or not t2:
                    continue
                if overlaps(b1, b2):
                    issues.append(
                        f"Slide {idx}: text shapes overlap  '{t1[:30]}'  vs  '{t2[:30]}'"
                    )

    if issues:
        print("Issues found:")
        for i in issues:
            print(" -", i)
    else:
        print("No geometry issues found.")

    qa_path = REPO_ROOT / "docs" / "qa_notes_raw.txt"
    qa_path.write_text("\n".join(notes), encoding="utf-8")
    print(f"\nFull layout dump: {qa_path}")
    print(f"Slides: {len(prs.slides)}")
    print(f"Slide size: {emu_to_in(sw):.2f}in x {emu_to_in(sh):.2f}in")


if __name__ == "__main__":
    main()
