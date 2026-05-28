"""Build the FemaleDiseasePrevalence.pptx presentation.

Originally specified to use pptxgenjs (Node). Node is unavailable in this
environment, so this Python equivalent uses python-pptx and produces the
same slide structure, palette and layout rules described in the brief.

Run:
    py scripts/build_presentation.py

Output:
    output/FemaleDiseasPrevalence.pptx
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_TICK_MARK
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Inches, Pt

REPO_ROOT = Path(__file__).resolve().parent.parent
DATA_PATH = REPO_ROOT / "data" / "disease_research.json"
OUTPUT_PATH = REPO_ROOT / "output" / "FemaleDiseasPrevalence.pptx"

PRIMARY = RGBColor(0x3D, 0x2B, 0x56)      # deep plum
SECONDARY = RGBColor(0x7B, 0x5E, 0xA7)    # soft violet (Female bars)
ACCENT = RGBColor(0xE8, 0xA8, 0x38)       # gold (Male bars / stat callouts)
SURFACE = RGBColor(0xF3, 0xEE, 0xF8)      # lavender mist (cards)
TEXT = RGBColor(0x2C, 0x2C, 0x2C)         # near-black body text
MUTED = RGBColor(0x88, 0x88, 0x88)        # footnotes
BG = RGBColor(0xFF, 0xF8, 0xF5)           # very light warm background
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BIG_IDEA = (
    "The same female-skewed health burden that shows up in Zambia's diabetes "
    "numbers runs through five of the most disabling diseases of modern life "
    "— and the world still designs medicine, research and screening as if it "
    "didn't."
)


def _set_background(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    font_size=18,
    bold=False,
    color=TEXT,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font_name=FONT,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _add_pill(slide, left, top, width, height, text, fill=SURFACE, font_color=PRIMARY, font_size=11):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.5
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(font_size)
    run.font.bold = True
    run.font.color.rgb = font_color
    return shape


def _add_card(slide, left, top, width, height, fill=SURFACE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.adjustments[0] = 0.06
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape


def _add_footnote(slide, text, color=MUTED):
    _add_text(
        slide,
        Inches(0.5),
        Inches(7.05),
        Inches(11.4),
        Inches(0.35),
        text,
        font_size=9,
        color=color,
        align=PP_ALIGN.LEFT,
    )


def _add_slide_number(slide, n, total):
    _add_text(
        slide,
        Inches(12.5),
        Inches(7.05),
        Inches(0.6),
        Inches(0.3),
        f"{n} / {total}",
        font_size=9,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def build_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide, BG)
    _add_text(
        slide,
        Inches(0.9),
        Inches(0.7),
        Inches(11.5),
        Inches(0.4),
        "A 3-MINUTE STORY",
        font_size=14,
        bold=True,
        color=SECONDARY,
    )
    _add_text(
        slide,
        Inches(0.9),
        Inches(1.4),
        Inches(11.5),
        Inches(1.8),
        "Diabetes was the question.\nThe answer is much bigger than diabetes.",
        font_size=44,
        bold=True,
        color=PRIMARY,
    )
    _add_text(
        slide,
        Inches(0.9),
        Inches(3.9),
        Inches(0.6),
        Inches(0.4),
        "Big Idea",
        font_size=11,
        bold=True,
        color=ACCENT,
    )
    _add_text(
        slide,
        Inches(0.9),
        Inches(4.35),
        Inches(11.5),
        Inches(2.2),
        BIG_IDEA,
        font_size=20,
        color=TEXT,
    )
    _add_text(
        slide,
        Inches(0.9),
        Inches(6.8),
        Inches(11.5),
        Inches(0.4),
        "Female disease prevalence  ·  Five conditions, one pattern",
        font_size=11,
        color=MUTED,
    )
    return slide


def build_hook_slide(prs, hook):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_text(
        slide,
        Inches(0.9),
        Inches(0.7),
        Inches(11.5),
        Inches(0.4),
        "THE HOOK  ·  Diabesity Zambia",
        font_size=12,
        bold=True,
        color=SECONDARY,
    )
    _add_text(
        slide,
        Inches(0.9),
        Inches(1.45),
        Inches(11.5),
        Inches(0.55),
        "Zambia's 2017 STEPS survey",
        font_size=22,
        color=PRIMARY,
        bold=True,
    )
    _add_text(
        slide,
        Inches(0.9),
        Inches(2.25),
        Inches(11.5),
        Inches(2.2),
        "24%",
        font_size=110,
        bold=True,
        color=ACCENT,
    )
    _add_text(
        slide,
        Inches(0.9),
        Inches(4.7),
        Inches(11.5),
        Inches(1.8),
        "more women than men carried multiple non-communicable disease "
        "risk factors — and female diabetes prevalence ran 3.0% vs 2.1% in men.",
        font_size=22,
        color=TEXT,
    )
    _add_footnote(
        slide,
        f"Source: {hook['citation']}  ·  {hook['source_url']}",
    )
    return slide


def build_landscape_slide(prs, diseases):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_text(
        slide,
        Inches(0.9),
        Inches(0.55),
        Inches(11.5),
        Inches(0.4),
        "THE LANDSCAPE",
        font_size=12,
        bold=True,
        color=SECONDARY,
    )
    _add_text(
        slide,
        Inches(0.9),
        Inches(1.0),
        Inches(11.5),
        Inches(0.7),
        "Female-to-male prevalence ratio, five conditions",
        font_size=26,
        bold=True,
        color=PRIMARY,
    )

    sorted_diseases = sorted(diseases, key=lambda d: d["female_to_male_ratio"])
    categories = [d["name"].split(" (")[0] for d in sorted_diseases]
    ratios = [d["female_to_male_ratio"] for d in sorted_diseases]

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Female-to-male ratio", ratios)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.9),
        Inches(1.95),
        Inches(11.5),
        Inches(4.9),
        chart_data,
    ).chart

    chart.has_title = False
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.font.size = Pt(16)
    dl.font.name = FONT
    dl.font.bold = True
    dl.font.color.rgb = PRIMARY
    dl.number_format = '0.0"x"'
    dl.position = XL_LABEL_POSITION.OUTSIDE_END

    series = plot.series[0]
    fill = series.format.fill
    fill.solid()
    fill.fore_color.rgb = SECONDARY
    series.format.line.fill.background()

    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(14)
    cat_axis.tick_labels.font.name = FONT
    cat_axis.tick_labels.font.color.rgb = TEXT

    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.major_tick_mark = XL_TICK_MARK.NONE
    val_axis.minor_tick_mark = XL_TICK_MARK.NONE
    val_axis.tick_labels.font.size = Pt(1)
    val_axis.maximum_scale = 10.0
    val_axis.minimum_scale = 0.0

    _add_text(
        slide,
        Inches(0.9),
        Inches(6.85),
        Inches(11.5),
        Inches(0.3),
        "Each bar = how many times more common the condition is in women than men.",
        font_size=11,
        color=MUTED,
        align=PP_ALIGN.LEFT,
    )
    return slide


def _add_stat_callout(slide, left, top, width, text, font_size=110, color=ACCENT, align=PP_ALIGN.LEFT):
    return _add_text(
        slide,
        left,
        top,
        width,
        Inches(2.0),
        text,
        font_size=font_size,
        bold=True,
        color=color,
        align=align,
    )


def build_disease_slide_layout_a_center_hero(prs, d, n, total):
    """Layout A — large stat centered, name above, so-what below. (Lupus)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_pill(slide, Inches(0.9), Inches(0.55), Inches(1.8), Inches(0.4), d["category"].upper())
    _add_text(
        slide,
        Inches(0.9),
        Inches(1.15),
        Inches(11.5),
        Inches(0.8),
        d["name"],
        font_size=34,
        bold=True,
        color=PRIMARY,
    )
    _add_text(
        slide,
        Inches(0.9),
        Inches(2.0),
        Inches(11.5),
        Inches(0.7),
        d["description"],
        font_size=16,
        color=TEXT,
    )
    _add_stat_callout(
        slide,
        Inches(0.9),
        Inches(3.2),
        Inches(11.5),
        d["ratio_label"],
        font_size=160,
        align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        Inches(0.9),
        Inches(5.6),
        Inches(11.5),
        Inches(0.6),
        "female-to-male prevalence",
        font_size=14,
        color=SECONDARY,
        align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        Inches(0.9),
        Inches(6.2),
        Inches(11.5),
        Inches(0.6),
        d["so_what"],
        font_size=20,
        bold=True,
        color=PRIMARY,
        align=PP_ALIGN.CENTER,
    )
    _add_footnote(slide, f"[{d['id']}] {d['citation']}  ·  {d['source_url']}")
    _add_slide_number(slide, n, total)
    return slide


def build_disease_slide_layout_b_two_column(prs, d, n, total):
    """Layout B — stat hero on the left, story column on the right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_pill(slide, Inches(0.9), Inches(0.55), Inches(2.0), Inches(0.4), d["category"].upper())

    # Left column — stat
    _add_text(
        slide,
        Inches(0.9),
        Inches(2.0),
        Inches(6.0),
        Inches(0.5),
        "female-to-male ratio",
        font_size=14,
        color=SECONDARY,
    )
    _add_stat_callout(
        slide,
        Inches(0.9),
        Inches(2.5),
        Inches(6.0),
        d["ratio_label"],
        font_size=170,
        align=PP_ALIGN.LEFT,
    )

    # Right column — story
    _add_text(
        slide,
        Inches(7.2),
        Inches(1.3),
        Inches(5.5),
        Inches(1.1),
        d["name"],
        font_size=30,
        bold=True,
        color=PRIMARY,
    )
    _add_text(
        slide,
        Inches(7.2),
        Inches(2.6),
        Inches(5.5),
        Inches(1.8),
        d["description"],
        font_size=16,
        color=TEXT,
    )
    _add_text(
        slide,
        Inches(7.2),
        Inches(4.8),
        Inches(5.5),
        Inches(1.8),
        d["so_what"],
        font_size=22,
        bold=True,
        color=PRIMARY,
    )
    _add_footnote(slide, f"[{d['id']}] {d['citation']}  ·  {d['source_url']}")
    _add_slide_number(slide, n, total)
    return slide


def build_disease_slide_layout_c_card_right(prs, d, n, total):
    """Layout C — story on the left, lavender card with stat on the right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_pill(slide, Inches(0.9), Inches(0.55), Inches(2.4), Inches(0.4), d["category"].upper())

    # Left — description
    _add_text(
        slide,
        Inches(0.9),
        Inches(1.3),
        Inches(6.2),
        Inches(1.1),
        d["name"],
        font_size=30,
        bold=True,
        color=PRIMARY,
    )
    _add_text(
        slide,
        Inches(0.9),
        Inches(2.6),
        Inches(6.2),
        Inches(2.5),
        d["description"],
        font_size=16,
        color=TEXT,
    )
    _add_text(
        slide,
        Inches(0.9),
        Inches(5.3),
        Inches(6.2),
        Inches(1.5),
        d["so_what"],
        font_size=20,
        bold=True,
        color=PRIMARY,
    )

    # Right — card with stat
    _add_card(slide, Inches(7.6), Inches(1.3), Inches(5.1), Inches(5.4), fill=SURFACE)
    _add_text(
        slide,
        Inches(7.6),
        Inches(1.7),
        Inches(5.1),
        Inches(0.5),
        "LIFETIME RISK GAP",
        font_size=12,
        bold=True,
        color=SECONDARY,
        align=PP_ALIGN.CENTER,
    )
    _add_stat_callout(
        slide,
        Inches(7.6),
        Inches(2.3),
        Inches(5.1),
        d["ratio_label"],
        font_size=130,
        align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        Inches(7.6),
        Inches(4.5),
        Inches(5.1),
        Inches(0.6),
        "female : male",
        font_size=16,
        color=SECONDARY,
        align=PP_ALIGN.CENTER,
    )
    _add_text(
        slide,
        Inches(7.6),
        Inches(5.2),
        Inches(5.1),
        Inches(1.4),
        "1 in 5 women  vs  1 in 10 men carry a lifetime Alzheimer's risk from age 45.",
        font_size=14,
        color=TEXT,
        align=PP_ALIGN.CENTER,
    )
    _add_footnote(slide, f"[{d['id']}] {d['citation']}  ·  {d['source_url']}")
    _add_slide_number(slide, n, total)
    return slide


def build_disease_slide_layout_d_top_stat(prs, d, n, total):
    """Layout D — stat at the top, name + story stacked below."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_pill(slide, Inches(0.9), Inches(0.55), Inches(2.0), Inches(0.4), d["category"].upper())

    _add_text(
        slide,
        Inches(0.9),
        Inches(1.15),
        Inches(11.5),
        Inches(0.35),
        "female-to-male ratio",
        font_size=14,
        color=SECONDARY,
        align=PP_ALIGN.LEFT,
    )
    _add_stat_callout(
        slide,
        Inches(0.9),
        Inches(1.6),
        Inches(11.5),
        d["ratio_label"],
        font_size=170,
        align=PP_ALIGN.LEFT,
    )

    _add_text(
        slide,
        Inches(0.9),
        Inches(4.4),
        Inches(11.5),
        Inches(0.9),
        d["name"],
        font_size=30,
        bold=True,
        color=PRIMARY,
    )
    _add_text(
        slide,
        Inches(0.9),
        Inches(5.3),
        Inches(11.5),
        Inches(0.9),
        d["description"],
        font_size=16,
        color=TEXT,
    )
    _add_text(
        slide,
        Inches(0.9),
        Inches(6.3),
        Inches(11.5),
        Inches(0.6),
        d["so_what"],
        font_size=18,
        bold=True,
        color=PRIMARY,
    )
    _add_footnote(slide, f"[{d['id']}] {d['citation']}  ·  {d['source_url']}")
    _add_slide_number(slide, n, total)
    return slide


def build_so_what_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_text(
        slide,
        Inches(0.9),
        Inches(0.7),
        Inches(11.5),
        Inches(0.4),
        "THE SO-WHAT",
        font_size=12,
        bold=True,
        color=SECONDARY,
    )
    _add_text(
        slide,
        Inches(0.9),
        Inches(1.2),
        Inches(0.6),
        Inches(0.4),
        "Big Idea",
        font_size=11,
        bold=True,
        color=ACCENT,
    )
    _add_text(
        slide,
        Inches(0.9),
        Inches(1.7),
        Inches(11.5),
        Inches(3.0),
        BIG_IDEA,
        font_size=26,
        bold=True,
        color=PRIMARY,
    )
    _add_card(slide, Inches(0.9), Inches(5.2), Inches(11.5), Inches(1.6), fill=SURFACE)
    _add_text(
        slide,
        Inches(1.2),
        Inches(5.4),
        Inches(11.0),
        Inches(0.4),
        "ONE RECOMMENDATION",
        font_size=11,
        bold=True,
        color=ACCENT,
    )
    _add_text(
        slide,
        Inches(1.2),
        Inches(5.85),
        Inches(11.0),
        Inches(1.0),
        "Make sex-disaggregated reporting the default in every NCD surveillance "
        "round — starting with Zambia's next STEPS — so the burden women carry "
        "stops hiding in the totals.",
        font_size=18,
        color=TEXT,
    )
    return slide


def build_sources_slide(prs, diseases, hook):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_background(slide)
    _add_text(
        slide,
        Inches(0.9),
        Inches(0.5),
        Inches(11.5),
        Inches(0.6),
        "Sources",
        font_size=26,
        bold=True,
        color=PRIMARY,
    )

    box = slide.shapes.add_textbox(Inches(0.9), Inches(1.3), Inches(11.5), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)

    entries = []
    for d in diseases:
        entries.append((f"[{d['id']}] {d['name']}", f"{d['citation']}.  {d['source_url']}"))
    entries.append(("[H] Hook — Zambia STEPS 2017", f"{hook['citation']}.  {hook['source_url']}"))

    for i, (head, body) in enumerate(entries):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run1 = p.add_run()
        run1.text = head + " — "
        run1.font.name = FONT
        run1.font.size = Pt(12)
        run1.font.bold = True
        run1.font.color.rgb = PRIMARY
        run2 = p.add_run()
        run2.text = body
        run2.font.name = FONT
        run2.font.size = Pt(11)
        run2.font.color.rgb = TEXT

    _add_footnote(
        slide,
        "All sources accessed via PubMed Central, CDC NCHS, and peer-reviewed journals.",
    )
    return slide


def main():
    data = json.loads(DATA_PATH.read_text(encoding="utf-8"))
    diseases = data["diseases"]
    hook = data["hook_anchor"]

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Map id -> layout function. Vary across the 5 disease slides (≥3 layouts).
    layout_for_id = {
        1: build_disease_slide_layout_a_center_hero,   # Lupus
        2: build_disease_slide_layout_b_two_column,    # Osteoporosis
        3: build_disease_slide_layout_d_top_stat,      # Migraine
        4: build_disease_slide_layout_c_card_right,    # Alzheimer's
        5: build_disease_slide_layout_b_two_column,    # Depression
    }

    total_slides = 10

    build_title_slide(prs)
    build_hook_slide(prs, hook)
    build_landscape_slide(prs, diseases)
    n = 4
    for d in diseases:
        layout_for_id[d["id"]](prs, d, n, total_slides)
        n += 1
    build_so_what_slide(prs)
    build_sources_slide(prs, diseases, hook)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    print(f"Wrote {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
